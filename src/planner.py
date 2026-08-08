from __future__ import annotations

from dataclasses import dataclass
from pathlib import Path
from typing import Iterable

from pptx import Presentation


@dataclass(frozen=True)
class SlideInfo:
    number: int
    text: str
    worktypes: frozenset[str]
    picture_count: int
    empty_picture_count: int


@dataclass(frozen=True)
class PlanStep:
    source: str
    slide: int
    section: str | None = None
    clean: bool = True
    worktype: str | None = None
    delete_empty_patrol_blocks: bool = False

    def as_dict(self) -> dict:
        result = {
            "source": self.source,
            "slide": self.slide,
            "clean": self.clean,
        }
        if self.section:
            result["section"] = self.section
        if self.worktype:
            result["worktype"] = self.worktype
        if self.delete_empty_patrol_blocks:
            result["delete_empty_patrol_blocks"] = True
        return result


def _slide_text(slide) -> str:
    chunks: list[str] = []
    for shape in slide.shapes:
        try:
            if hasattr(shape, "text") and shape.text.strip():
                chunks.append(shape.text.strip())
        except Exception:
            pass
        try:
            if shape.has_table:
                for row in shape.table.rows:
                    chunks.append(" | ".join(cell.text.strip() for cell in row.cells))
        except Exception:
            pass
    return "\n".join(chunks)


def _table_worktypes(slide) -> set[str]:
    result: set[str] = set()
    for shape in slide.shapes:
        try:
            if not shape.has_table:
                continue
            table = shape.table
            header_row = None
            work_col = None
            for row_index in range(min(3, len(table.rows))):
                headers = [
                    "".join(ch.lower() for ch in cell.text if ch.isalnum())
                    for cell in table.rows[row_index].cells
                ]
                for col_index, header in enumerate(headers):
                    if "worktype" in header:
                        header_row = row_index
                        work_col = col_index
                        break
                if work_col is not None:
                    break
            if work_col is None:
                continue
            for row_index in range((header_row or 0) + 1, len(table.rows)):
                value = table.rows[row_index].cells[work_col].text.strip().upper()
                if value:
                    result.add(value)
        except Exception:
            pass
    return result


def inspect_presentation(path: Path) -> list[SlideInfo]:
    prs = Presentation(path)
    infos: list[SlideInfo] = []
    for number, slide in enumerate(prs.slides, start=1):
        pictures = 0
        empty = 0
        for shape in slide.shapes:
            try:
                if int(shape.shape_type) not in {13, 14}:
                    continue
                blip = shape.element.find(
                    ".//{http://schemas.openxmlformats.org/drawingml/2006/main}blip"
                )
                rel_id = None
                if blip is not None:
                    rel_id = blip.get(
                        "{http://schemas.openxmlformats.org/officeDocument/2006/relationships}embed"
                    )
                if rel_id:
                    pictures += 1
                elif int(shape.shape_type) == 14:
                    empty += 1
            except Exception:
                pass
        infos.append(
            SlideInfo(
                number=number,
                text=_slide_text(slide),
                worktypes=frozenset(_table_worktypes(slide)),
                picture_count=pictures,
                empty_picture_count=empty,
            )
        )
    return infos


def _contains(info: SlideInfo, *phrases: str) -> bool:
    lower = info.text.lower()
    return all(phrase.lower() in lower for phrase in phrases)


def _first(infos: Iterable[SlideInfo], predicate, description: str) -> SlideInfo:
    for info in infos:
        if predicate(info):
            return info
    raise RuntimeError(f"Could not locate slide: {description}")


def _all(infos: Iterable[SlideInfo], predicate) -> list[SlideInfo]:
    return [info for info in infos if predicate(info)]


def _section_slides(
    infos: list[SlideInfo],
    title_prefixes: tuple[str, ...],
    worktype: str,
) -> list[SlideInfo]:
    result: list[SlideInfo] = []
    for info in infos:
        lower = info.text.lower()
        title_match = any(prefix in lower for prefix in title_prefixes)
        if title_match or worktype.upper() in info.worktypes:
            result.append(info)
    return result


def build_slide_plan(sources: dict[str, Path]) -> list[dict]:
    """
    Build the report order from slide content, not fixed slide numbers.

    This tolerates inserted/deleted slides and moderate formatting changes. The
    report's business section order remains fixed, while each source slide is
    discovered from titles, table work types, and section boundaries.
    """
    info = {key: inspect_presentation(path) for key, path in sources.items()}
    plan: list[PlanStep] = []

    # Intro and weather.
    cover = _first(
        info["CHEC"],
        lambda s: _contains(s, "patrol diary", "roads maintenance department"),
        "CHEC cover",
    )
    project = _first(info["CHEC"], lambda s: _contains(s, "project title", "work order number"), "project details")
    location = _first(info["CHEC"], lambda s: _contains(s, "location map"), "location map")
    weather = None
    for source_key in ("ALCAT", "SATURN", "CHEC", "TTG"):
        matches = _all(info[source_key], lambda s: _contains(s, "weather report"))
        if matches:
            weather = (source_key, matches[0])
            break
    if weather is None:
        raise RuntimeError("Could not locate a Weather Report slide.")

    plan.extend(
        [
            PlanStep("CHEC", cover.number, clean=False),
            PlanStep("CHEC", project.number, clean=False),
            PlanStep("CHEC", location.number, clean=False),
            PlanStep(weather[0], weather[1].number, clean=False),
        ]
    )

    # Road coverage / ODO slides, grouped by subcontractor.
    for source_key in ("CHEC", "ALCAT", "SATURN"):
        coverage = _all(
            info[source_key],
            lambda s: "length covered" in s.text.lower()
            and "route reference" in s.text.lower()
            and not s.worktypes,
        )
        if not coverage:
            raise RuntimeError(f"No road-coverage slides found in {source_key}.")
        for slide in coverage:
            plan.append(
                PlanStep(
                    source_key,
                    slide.number,
                    clean=True,
                    delete_empty_patrol_blocks=True,
                )
            )

    # Patrolling tables: exactly one section slide per source.
    for source_key in ("CHEC", "ALCAT", "SATURN", "TTG"):
        candidates = _section_slides(
            info[source_key],
            ("patrolling during the reporting date", "maximo during the reporting date"),
            "PS",
        )
        slide = _first(
            candidates,
            lambda s: "ps" in s.worktypes
            or "patrolling during the reporting date" in s.text.lower()
            or "maximo during the reporting date" in s.text.lower(),
            f"{source_key} patrolling section",
        )
        plan.append(PlanStep(source_key, slide.number, section="PS", worktype="PS"))

    # CM, FR, and IN sections.
    section_rules = [
        ("CM", ("corrective maintenance during the reporting date", "details of rectified defects during the reporting date-cm"), ("CHEC", "ALCAT", "SATURN", "TTG")),
        ("FR", ("first response during the reporting date", "details of rectified defects during the reporting date-fr"), ("CHEC", "ALCAT", "SATURN", "TTG")),
        ("IN", ("in during the reporting date", "inspection during the reporting date", "details of rectified defects during the reporting date-in"), ("CHEC", "ALCAT", "SATURN")),
    ]
    last_section_slide: dict[str, int] = {}
    for worktype, prefixes, source_order in section_rules:
        for source_key in source_order:
            slides = _section_slides(info[source_key], prefixes, worktype)
            # The CHEC Maximo overview contains all worktypes and must not be
            # repeated in CM/FR/IN; dedicated section titles take precedence.
            dedicated = [
                s for s in slides
                if any(prefix in s.text.lower() for prefix in prefixes)
            ]
            if dedicated:
                slides = dedicated
            if not slides:
                continue
            for slide in slides:
                plan.append(
                    PlanStep(
                        source_key,
                        slide.number,
                        section=worktype,
                        worktype=worktype,
                    )
                )
                last_section_slide[source_key] = max(last_section_slide.get(source_key, 0), slide.number)

    # Progress/asset slides are everything after the final operational section
    # in each source, excluding known summary-only pages without photos/assets.
    for source_key in ("CHEC", "ALCAT", "SATURN", "TTG"):
        start_after = last_section_slide.get(source_key, 0)
        trailing = [s for s in info[source_key] if s.number > start_after]
        if source_key == "CHEC":
            trailing = [
                s for s in trailing
                if s.picture_count > 0
                or "progress photos" in s.text.lower()
                or "wo no" in s.text.lower()
            ]
        elif source_key == "ALCAT":
            trailing = [
                s for s in trailing
                if not _contains(s, "number of work orders created by inspectors per day")
            ]
        for slide in trailing:
            plan.append(PlanStep(source_key, slide.number, clean=True))

    if len(plan) < 40:
        raise RuntimeError(
            f"Dynamic planner found only {len(plan)} output slides; expected a full patrol diary."
        )
    return [step.as_dict() for step in plan]
