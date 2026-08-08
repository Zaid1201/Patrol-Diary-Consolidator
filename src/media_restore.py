from __future__ import annotations

import hashlib
import io
import re
from dataclasses import dataclass
from functools import lru_cache
from pathlib import Path
from typing import Any, Iterable

from PIL import Image, ImageOps
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE

# Office constants used through COM.
msoFalse = 0
msoTrue = -1
msoSendToBack = 1
msoBringForward = 2
msoBringToFront = 0
msoPicture = 13
msoLinkedPicture = 11
msoPlaceholder = 14
ppPlaceholderPicture = 18

EMU_PER_POINT = 12700.0

A_NS = "http://schemas.openxmlformats.org/drawingml/2006/main"
R_NS = "http://schemas.openxmlformats.org/officeDocument/2006/relationships"
P_NS = "http://schemas.openxmlformats.org/presentationml/2006/main"
NS = {"a": A_NS, "r": R_NS, "p": P_NS}


@dataclass(frozen=True)
class Rect:
    left: float
    top: float
    width: float
    height: float

    @property
    def right(self) -> float:
        return self.left + self.width

    @property
    def bottom(self) -> float:
        return self.top + self.height

    @property
    def center_x(self) -> float:
        return self.left + self.width / 2.0

    @property
    def center_y(self) -> float:
        return self.top + self.height / 2.0


@dataclass(frozen=True)
class ImagePlacement:
    file_path: Path
    rect: Rect
    source_z_index: int
    rotation: float
    digest: str


@dataclass(frozen=True)
class EmptyPictureSlot:
    rect: Rect
    source_z_index: int


@dataclass(frozen=True)
class SlideMedia:
    images: tuple[ImagePlacement, ...]
    empty_slots: tuple[EmptyPictureSlot, ...]
    slide_width: float
    slide_height: float


def emu_to_points(value: int | float) -> float:
    return float(value) / EMU_PER_POINT


@lru_cache(maxsize=16)
def _load_presentation(path_text: str) -> Presentation:
    return Presentation(path_text)


def _local_name(element: Any) -> str:
    tag = getattr(element, "tag", "")
    return tag.split("}")[-1]


def _is_picture_placeholder(shape: Any) -> bool:
    try:
        return (
            shape.shape_type == MSO_SHAPE_TYPE.PLACEHOLDER
            and int(shape.placeholder_format.type) == ppPlaceholderPicture
        )
    except Exception:
        name = str(getattr(shape, "name", "")).lower()
        return "picture placeholder" in name


def _image_extension(part: Any) -> str:
    suffix = Path(str(getattr(part, "partname", ""))).suffix.lower().lstrip(".")
    if suffix == "jpeg":
        suffix = "jpg"
    if suffix in {"png", "jpg", "gif", "bmp", "tif", "tiff", "emf", "wmf"}:
        return suffix

    content_type = str(getattr(part, "content_type", "")).lower()
    mapping = {
        "image/png": "png",
        "image/jpeg": "jpg",
        "image/gif": "gif",
        "image/bmp": "bmp",
        "image/tiff": "tif",
        "image/x-emf": "emf",
        "image/x-wmf": "wmf",
    }
    return mapping.get(content_type, "png")


def _crop_values(shape: Any) -> tuple[int, int, int, int]:
    src_rect = shape.element.find(".//a:srcRect", NS)
    if src_rect is None:
        return 0, 0, 0, 0

    def val(name: str) -> int:
        try:
            return max(0, min(100000, int(src_rect.get(name, "0"))))
        except Exception:
            return 0

    return val("l"), val("t"), val("r"), val("b")


def _write_image(blob: bytes, extension: str, crop: tuple[int, int, int, int], output_base: Path) -> Path:
    """Write an Office image to disk, applying the PPT crop when present."""
    left, top, right, bottom = crop
    has_crop = any(crop)

    # Preserve the original bytes when no crop is required. This is faster and
    # avoids recompressing high-resolution site photographs.
    if not has_crop and extension in {"png", "jpg", "gif", "bmp", "tif", "tiff", "emf", "wmf"}:
        path = output_base.with_suffix(f".{extension}")
        path.write_bytes(blob)
        return path

    try:
        with Image.open(io.BytesIO(blob)) as opened:
            image = ImageOps.exif_transpose(opened)
            image.load()
            width, height = image.size

            x1 = int(round(width * left / 100000.0))
            y1 = int(round(height * top / 100000.0))
            x2 = int(round(width * (1.0 - right / 100000.0)))
            y2 = int(round(height * (1.0 - bottom / 100000.0)))

            x1 = max(0, min(width - 1, x1))
            y1 = max(0, min(height - 1, y1))
            x2 = max(x1 + 1, min(width, x2))
            y2 = max(y1 + 1, min(height, y2))

            image = image.crop((x1, y1, x2, y2))
            if image.mode not in {"RGB", "RGBA"}:
                image = image.convert("RGBA" if "A" in image.getbands() else "RGB")

            path = output_base.with_suffix(".png")
            image.save(path, format="PNG", optimize=False)
            return path
    except Exception:
        # PowerPoint can read EMF/WMF and several formats that Pillow cannot.
        path = output_base.with_suffix(f".{extension}")
        path.write_bytes(blob)
        return path


def extract_slide_media(source_path: Path, slide_number: int, temp_root: Path) -> SlideMedia:
    source_path = source_path.resolve()
    prs = _load_presentation(str(source_path))
    if slide_number < 1 or slide_number > len(prs.slides):
        raise IndexError(f"Slide {slide_number} does not exist in {source_path.name}")

    slide = prs.slides[slide_number - 1]
    slide_folder = temp_root / re.sub(r"[^A-Za-z0-9_.-]+", "_", source_path.stem) / f"slide_{slide_number:03d}"
    slide_folder.mkdir(parents=True, exist_ok=True)

    images: list[ImagePlacement] = []
    empty_slots: list[EmptyPictureSlot] = []

    for z_index, shape in enumerate(slide.shapes, start=1):
        element_name = _local_name(shape.element)
        is_picture_element = element_name == "pic"
        is_picture_placeholder = _is_picture_placeholder(shape)

        if not (is_picture_element or is_picture_placeholder):
            continue

        rect = Rect(
            emu_to_points(shape.left),
            emu_to_points(shape.top),
            emu_to_points(shape.width),
            emu_to_points(shape.height),
        )

        blip = shape.element.find(".//a:blip", NS)
        rel_id = blip.get(f"{{{R_NS}}}embed") if blip is not None else None

        if not rel_id:
            empty_slots.append(EmptyPictureSlot(rect=rect, source_z_index=z_index))
            continue

        try:
            image_part = slide.part.related_part(rel_id)
            blob = image_part.blob
        except Exception:
            empty_slots.append(EmptyPictureSlot(rect=rect, source_z_index=z_index))
            continue

        digest = hashlib.sha256(blob).hexdigest()
        extension = _image_extension(image_part)
        crop = _crop_values(shape)
        output_base = slide_folder / f"z{z_index:03d}_{digest[:12]}"
        image_path = _write_image(blob, extension, crop, output_base)

        try:
            rotation = float(shape.rotation or 0.0)
        except Exception:
            rotation = 0.0

        images.append(
            ImagePlacement(
                file_path=image_path,
                rect=rect,
                source_z_index=z_index,
                rotation=rotation,
                digest=digest,
            )
        )

    return SlideMedia(
        images=tuple(images),
        empty_slots=tuple(empty_slots),
        slide_width=emu_to_points(prs.slide_width),
        slide_height=emu_to_points(prs.slide_height),
    )


def extract_standard_logos(source_path: Path, temp_root: Path, slide_number: int = 2) -> tuple[ImagePlacement, ...]:
    media = extract_slide_media(source_path, slide_number, temp_root)
    logos: list[ImagePlacement] = []

    # The standard report header uses four image logos within the top 100 pt,
    # split between the left and right corners.
    for image in media.images:
        in_top_band = image.rect.top < 105 and image.rect.bottom < 115
        in_corner = image.rect.left < 245 or image.rect.right > media.slide_width - 245
        if in_top_band and in_corner:
            logos.append(image)

    if len(logos) < 2:
        raise RuntimeError(
            f"Could not find the standard header logos in {source_path.name} slide {slide_number}."
        )
    return tuple(logos)


def _com_is_picture_like(shape: Any) -> bool:
    try:
        shape_type = int(shape.Type)
    except Exception:
        return False

    if shape_type in {msoPicture, msoLinkedPicture}:
        return True

    if shape_type == msoPlaceholder:
        try:
            return int(shape.PlaceholderFormat.Type) == ppPlaceholderPicture
        except Exception:
            name = str(getattr(shape, "Name", "")).lower()
            return "picture" in name
    return False


def delete_slide_picture_shapes(slide: Any) -> int:
    deleted = 0
    for index in range(slide.Shapes.Count, 0, -1):
        shape = slide.Shapes.Item(index)
        if _com_is_picture_like(shape):
            try:
                shape.Delete()
                deleted += 1
            except Exception:
                pass
    return deleted


def _set_image_z_order(shape: Any, target_index: int):
    """Place an added image at its original 1-based z-order index."""
    try:
        shape.ZOrder(msoSendToBack)
        for _ in range(max(0, target_index - 1)):
            shape.ZOrder(msoBringForward)
    except Exception:
        pass


def _add_picture(slide: Any, image: ImagePlacement, z_index: int | None = None) -> Any:
    shape = slide.Shapes.AddPicture(
        str(image.file_path.resolve()),
        msoFalse,
        msoTrue,
        image.rect.left,
        image.rect.top,
        image.rect.width,
        image.rect.height,
    )
    try:
        shape.Rotation = image.rotation
    except Exception:
        pass

    if z_index is not None:
        _set_image_z_order(shape, z_index)
    return shape


def _horizontal_overlap(a: Rect, b: Rect) -> float:
    return max(0.0, min(a.right, b.right) - max(a.left, b.left))


def _merge_empty_columns(empty_slots: Iterable[EmptyPictureSlot], max_gap: float = 25.0) -> list[Rect]:
    slots = sorted((slot.rect for slot in empty_slots), key=lambda r: (r.left, r.top))
    groups: list[Rect] = []

    for rect in slots:
        merged = False
        for i, group in enumerate(groups):
            overlap = _horizontal_overlap(group, rect)
            gap = max(0.0, max(group.left, rect.left) - min(group.right, rect.right))
            if overlap > 0 or gap <= max_gap:
                left = min(group.left, rect.left)
                right = max(group.right, rect.right)
                top = min(group.top, rect.top)
                bottom = max(group.bottom, rect.bottom)
                groups[i] = Rect(left, top, right - left, bottom - top)
                merged = True
                break
        if not merged:
            groups.append(rect)

    # A merge can make two earlier groups adjacent. Collapse until stable.
    changed = True
    while changed:
        changed = False
        result: list[Rect] = []
        for rect in groups:
            for i, group in enumerate(result):
                overlap = _horizontal_overlap(group, rect)
                gap = max(0.0, max(group.left, rect.left) - min(group.right, rect.right))
                if overlap > 0 or gap <= max_gap:
                    left = min(group.left, rect.left)
                    right = max(group.right, rect.right)
                    top = min(group.top, rect.top)
                    bottom = max(group.bottom, rect.bottom)
                    result[i] = Rect(left, top, right - left, bottom - top)
                    changed = True
                    break
            else:
                result.append(rect)
        groups = result
    return groups


def _shape_text(shape: Any) -> str:
    try:
        if shape.HasTextFrame and shape.TextFrame.HasText:
            return str(shape.TextFrame.TextRange.Text or "")
    except Exception:
        pass
    return ""


def _shape_rect(shape: Any) -> Rect:
    return Rect(float(shape.Left), float(shape.Top), float(shape.Width), float(shape.Height))


def _is_empty_description(text: str) -> bool:
    normalized = re.sub(r"[^a-z0-9]+", " ", text.lower()).strip()
    normalized = normalized.replace("wo no", "").replace("description", "").strip()
    return normalized == ""


def _vertical_overlap(a: Rect, b: Rect) -> float:
    return max(0.0, min(a.bottom, b.bottom) - max(a.top, b.top))


def _rect_matches_slot(rect: Rect, slot: Rect, tolerance: float = 12.0) -> bool:
    center_match = (
        abs(rect.center_x - slot.center_x) <= max(tolerance, slot.width * 0.08)
        and abs(rect.center_y - slot.center_y) <= max(tolerance, slot.height * 0.08)
    )
    size_match = (
        abs(rect.width - slot.width) <= max(tolerance, slot.width * 0.10)
        and abs(rect.height - slot.height) <= max(tolerance, slot.height * 0.10)
    )
    return center_match and size_match


def cleanup_empty_picture_slots(slide: Any, media: SlideMedia) -> int:
    """
    Remove genuinely empty source picture slots and only their associated text.

    The source PPTX package is the authority: a placeholder is empty only when
    it has no embedded image relationship. Actual photographs are never moved or
    rebuilt, so Before/After work-order pairing remains unchanged.
    """
    deleted = 0
    body_slots = [slot for slot in media.empty_slots if slot.rect.top >= 105]

    # Always remove blank work-order caption boxes such as "WO No.: / Description:".
    # These boxes belong to unused progress-photo columns.
    if not body_slots:
        for index in range(slide.Shapes.Count, 0, -1):
            shape = slide.Shapes.Item(index)
            text = _shape_text(shape)
            if text and _is_empty_description(text):
                try:
                    shape.Delete()
                    deleted += 1
                except Exception:
                    pass
        return deleted

    occupied = [image.rect for image in media.images if image.rect.top >= 105]
    empty_groups = _merge_empty_columns(body_slots)

    # A whole column is empty only if no real image overlaps it both
    # horizontally and vertically. Header logos therefore cannot make a blank
    # body column look occupied.
    full_empty_regions: list[Rect] = []
    for group in empty_groups:
        has_picture = any(
            _horizontal_overlap(group, picture) >= min(group.width, picture.width) * 0.35
            and _vertical_overlap(group, picture) >= min(group.height, picture.height) * 0.20
            for picture in occupied
        )
        if not has_picture:
            full_empty_regions.append(group)

    for index in range(slide.Shapes.Count, 0, -1):
        shape = slide.Shapes.Item(index)
        try:
            rect = _shape_rect(shape)
        except Exception:
            continue
        text = _shape_text(shape)
        delete = bool(text and _is_empty_description(text))

        # Delete the copied placeholder itself when its geometry matches a
        # source placeholder that has no embedded image.
        if not delete:
            try:
                shape_type = int(shape.Type)
                shape_name = str(getattr(shape, "Name", "")).lower()
            except Exception:
                shape_type = 0
                shape_name = ""
            is_empty_placeholder_shape = (
                shape_type == msoPlaceholder or "placeholder" in shape_name
            )
            if is_empty_placeholder_shape:
                for slot in body_slots:
                    if _rect_matches_slot(rect, slot.rect):
                        delete = True
                        break

        # Remove Before/After labels local to an empty slot.
        compact = re.sub(r"\s+", "", text).lower()
        if not delete and compact in {"before:", "after:"}:
            for slot in body_slots:
                margin = 10.0
                if (
                    slot.rect.left - margin <= rect.center_x <= slot.rect.right + margin
                    and slot.rect.top - margin <= rect.center_y <= slot.rect.bottom + margin
                ):
                    delete = True
                    break

        # If both Before and After are empty, remove the entire work-order/road
        # column below the header, including borders and caption text.
        if not delete and rect.top > 105 and rect.width < media.slide_width * 0.72:
            for region in full_empty_regions:
                overlap = _horizontal_overlap(rect, region)
                center_in_region = region.left - 8 <= rect.center_x <= region.right + 8
                substantial = overlap >= min(rect.width, region.width) * 0.45
                if center_in_region and substantial:
                    delete = True
                    break

        if delete:
            try:
                shape.Delete()
                deleted += 1
            except Exception:
                pass

    return deleted

def restore_slide_images(
    slide: Any,
    source_path: Path,
    source_slide_number: int,
    temp_root: Path,
) -> dict[str, int]:
    """Replace copied/broken picture placeholders with actual embedded images."""
    media = extract_slide_media(source_path, source_slide_number, temp_root)
    removed = delete_slide_picture_shapes(slide)

    # Add in source order and restore the original z-order index. This keeps
    # photographs behind the Before/After labels but above background objects.
    added = 0
    for image in media.images:
        _add_picture(slide, image, z_index=image.source_z_index)
        added += 1

    cleaned = cleanup_empty_picture_slots(slide, media)
    return {"removed": removed, "added": added, "empty_cleaned": cleaned}


def ensure_standard_logos(
    slide: Any,
    logos: tuple[ImagePlacement, ...],
    slide_width: float,
) -> int:
    """Overlay the standard report logos without deleting source media.

    Native slide copy normally carries the source master correctly. Overlaying
    the same standard logos is harmless when they already exist and guarantees
    that slides whose source master omits them still match the consolidated
    report. No body picture is ever deleted by this function.
    """
    added = 0
    for logo in logos:
        shape = _add_picture(slide, logo, z_index=None)
        try:
            shape.ZOrder(msoBringToFront)
        except Exception:
            pass
        added += 1
    return added

